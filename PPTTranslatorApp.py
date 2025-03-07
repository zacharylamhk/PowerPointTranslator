import tkinter as tk
from tkinter import filedialog, ttk, messagebox
from pptx import Presentation
from deep_translator import GoogleTranslator
import os
from pptx.util import Pt
import sys
from io import StringIO

class TextRedirector:
    """Redirect stdout to a Tkinter Text widget"""
    def __init__(self, widget):
        self.widget = widget
        self.buffer = StringIO()

    def write(self, string):
        self.buffer.write(string)
        self.widget.insert(tk.END, string)
        self.widget.see(tk.END)
        self.widget.update_idletasks()

    def flush(self):
        self.buffer.seek(0)
        self.buffer.truncate(0)

def translate_text(text, source_language, target_language):
    """Translate text from source language to target language"""
    try:
        if not isinstance(text, str) or not text.strip():
            print(f"Skipping invalid text: {text}")
            return text if isinstance(text, str) else ""
        translator = GoogleTranslator(source=source_language, target=target_language)
        translated = translator.translate(text)
        if translated is None:
            print(f"Translation returned None for '{text}', keeping original")
            return text
        print(f"Translated '{text}' to '{translated}'")
        return translated
    except Exception as e:
        print(f"Translation error for text '{text}': {e}")
        return text if isinstance(text, str) else ""

def get_font_properties(run):
    """Extract font properties from a run"""
    font = run.font
    props = {
        'size': font.size.pt if font.size else None,
        'name': font.name if font.name else None,
        'bold': font.bold if font.bold is not None else False,
        'italic': font.italic if font.italic is not None else False,
        'underline': font.underline if font.underline is not None else False
    }
    print(f"Font properties for '{run.text}': {props}")
    return props

def apply_font_properties(run, properties):
    """Apply font properties to a run with detailed error catching"""
    font = run.font
    try:
        # Log before applying each property
        print(f"Applying properties to '{run.text}': {properties}")
        
        if properties['size'] is not None:
            try:
                font.size = Pt(properties['size'])
                print(f"Set size to {properties['size']} for '{run.text}'")
            except AttributeError as e:
                print(f"Failed to set size for '{run.text}': {e}")
        
        # Temporarily skip name to isolate the issue
        if properties['name'] is not None:
            print(f"Skipping font name '{properties['name']}' for '{run.text}' to avoid error")
            # Uncomment below to test if name is the issue
            # try:
            #     font.name = properties['name']
            #     print(f"Set name to {properties['name']} for '{run.text}'")
            # except AttributeError as e:
            #     print(f"Failed to set name for '{run.text}': {e}")
        
        try:
            font.bold = properties['bold'] if properties['bold'] is not None else False
            print(f"Set bold to {properties['bold']} for '{run.text}'")
        except AttributeError as e:
            print(f"Failed to set bold for '{run.text}': {e}")
        
        try:
            font.italic = properties['italic'] if properties['italic'] is not None else False
            print(f"Set italic to {properties['italic']} for '{run.text}'")
        except AttributeError as e:
            print(f"Failed to set italic for '{run.text}': {e}")
        
        try:
            font.underline = properties['underline'] if properties['underline'] is not None else False
            print(f"Set underline to {properties['underline']} for '{run.text}'")
        except AttributeError as e:
            print(f"Failed to set underline for '{run.text}': {e}")
            
    except Exception as e:
        print(f"Unexpected error applying font properties to '{run.text}': {e}")

def translate_ppt(input_path, output_path, source_lang, target_lang, progress_label, total_slides, log_widget):
    """Translate PowerPoint file with page number updates and preserve font"""
    try:
        prs = Presentation(input_path)
        slide_count = len(prs.slides)
        
        for i, slide in enumerate(prs.slides, 1):
            for shape in slide.shapes:
                if hasattr(shape, "text_frame") and shape.text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            if run.text:
                                font_props = get_font_properties(run)
                                translated = translate_text(run.text, source_lang, target_lang)
                                run.text = translated if translated is not None else run.text
                                apply_font_properties(run, font_props)
                
                if shape.has_table:
                    for row in shape.table.rows:
                        for cell in row.cells:
                            for paragraph in cell.text_frame.paragraphs:
                                for run in paragraph.runs:
                                    if run.text:
                                        font_props = get_font_properties(run)
                                        translated = translate_text(run.text, source_lang, target_lang)
                                        run.text = translated if translated is not None else run.text
                                        apply_font_properties(run, font_props)
                
                if shape.has_chart:
                    chart = shape.chart
                    for series in chart.series:
                        if series.name:
                            translated = translate_text(series.name, source_lang, target_lang)
                            series.name = translated if translated is not None else series.name
                    if chart.has_title:
                        for paragraph in chart.chart_title.text_frame.paragraphs:
                            for run in paragraph.runs:
                                if run.text:
                                    font_props = get_font_properties(run)
                                    translated = translate_text(run.text, source_lang, target_lang)
                                    run.text = translated if translated is not None else run.text
                                    apply_font_properties(run, font_props)
            
            progress_label.config(text=f"{i}/{slide_count}")
            progress_label.update()

        prs.save(output_path)
        messagebox.showinfo("Success", f"Translation completed!\nSaved as: {output_path}")
        
    except Exception as e:
        messagebox.showerror("Error", f"Error processing presentation: {e}")
    finally:
        progress_label.config(text="0/0")

class PPTTranslatorApp:
    def __init__(self, root):
        self.root = root
        self.root.title("PowerPoint Translator")
        self.root.geometry("500x400")

        # Variables
        self.input_path = tk.StringVar()
        self.output_path = tk.StringVar()
        self.source_lang = tk.StringVar(value='auto')
        self.target_lang = tk.StringVar(value='en')

        # Get supported languages
        self.languages = GoogleTranslator().get_supported_languages(as_dict=True)

        # GUI Elements
        tk.Label(root, text="Source PPTX:").grid(row=0, column=0, padx=5, pady=5, sticky="e")
        tk.Entry(root, textvariable=self.input_path, width=40).grid(row=0, column=1, padx=5, pady=5)
        tk.Button(root, text="Browse", command=self.browse_input).grid(row=0, column=2, padx=5, pady=5)

        tk.Label(root, text="Output PPTX:").grid(row=1, column=0, padx=5, pady=5, sticky="e")
        tk.Entry(root, textvariable=self.output_path, width=40).grid(row=1, column=1, padx=5, pady=5)
        tk.Button(root, text="Browse", command=self.browse_output).grid(row=1, column=2, padx=5, pady=5)

        tk.Label(root, text="Source Language:").grid(row=2, column=0, padx=5, pady=5, sticky="e")
        source_combo = ttk.Combobox(root, textvariable=self.source_lang, width=37)
        source_combo['values'] = ['auto'] + list(self.languages.keys())
        source_combo.grid(row=2, column=1, padx=5, pady=5)

        tk.Label(root, text="Target Language:").grid(row=3, column=0, padx=5, pady=5, sticky="e")
        target_combo = ttk.Combobox(root, textvariable=self.target_lang, width=37)
        target_combo['values'] = list(self.languages.keys())
        target_combo.grid(row=3, column=1, padx=5, pady=5)

        tk.Label(root, text="Progress:").grid(row=4, column=0, padx=5, pady=5, sticky="e")
        self.progress_label = tk.Label(root, text="0/0", width=40)
        self.progress_label.grid(row=4, column=1, columnspan=2, padx=5, pady=5)

        tk.Label(root, text="Log:").grid(row=5, column=0, padx=5, pady=5, sticky="ne")
        self.log_text = tk.Text(root, height=8, width=40, state='normal')
        self.log_text.grid(row=5, column=1, columnspan=2, padx=5, pady=5)
        self.log_text.config(state='disabled')

        tk.Button(root, text="Translate", command=self.translate).grid(row=6, column=1, pady=20)

    def browse_input(self):
        """Open file dialog for input file"""
        filename = filedialog.askopenfilename(
            filetypes=[("PowerPoint files", "*.pptx *.ppt"), ("All files", "*.*")]
        )
        if filename:
            self.input_path.set(filename)
            output_suggestion = f"translated_{os.path.basename(filename)}"
            self.output_path.set(os.path.join(os.path.dirname(filename), output_suggestion))

    def browse_output(self):
        """Open file dialog for output file"""
        filename = filedialog.asksaveasfilename(
            defaultextension=".pptx",
            filetypes=[("PowerPoint files", "*.pptx"), ("All files", "*.*")]
        )
        if filename:
            self.output_path.set(filename)

    def translate(self):
        """Handle translation process"""
        if not self.input_path.get():
            messagebox.showerror("Error", "Please select an input file!")
            return
        if not self.output_path.get():
            messagebox.showerror("Error", "Please specify an output file!")
            return
        if not os.path.exists(self.input_path.get()):
            messagebox.showerror("Error", "Input file not found!")
            return
        
        self.root.children['!button3'].config(state='disabled')
        self.log_text.config(state='normal')
        self.log_text.delete(1.0, tk.END)
        
        original_stdout = sys.stdout
        sys.stdout = TextRedirector(self.log_text)
        
        try:
            prs = Presentation(self.input_path.get())
            total_slides = len(prs.slides)
            translate_ppt(
                self.input_path.get(),
                self.output_path.get(),
                self.source_lang.get(),
                self.target_lang.get(),
                self.progress_label,
                total_slides,
                self.log_text
            )
        finally:
            sys.stdout = original_stdout
            self.log_text.config(state='disabled')
            self.root.children['!button3'].config(state='normal')

def main():
    root = tk.Tk()
    app = PPTTranslatorApp(root)
    root.mainloop()

if __name__ == "__main__":
    main()